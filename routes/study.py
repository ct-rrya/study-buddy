"""
Study Routes - File upload, Quiz, Bot interaction
"""
from flask import Blueprint, render_template, request, jsonify
from flask_login import login_required, current_user
from werkzeug.utils import secure_filename
from datetime import datetime
from models import db, StudyFile, StudySession, BotConversation, Subject, SubjectProgress
from services.bot import StudyBot
import os

study_bp = Blueprint('study', __name__)
ALLOWED_EXTENSIONS = {'txt', 'md', 'pptx', 'docx', 'xlsx', 'pdf'}

def extract_text_from_file(file, filename):
    """Extract text content from various file types"""
    ext = filename.rsplit('.', 1)[1].lower()
    
    if ext in ['txt', 'md']:
        return file.read().decode('utf-8')
    
    elif ext == 'pptx':
        from pptx import Presentation
        from io import BytesIO
        prs = Presentation(BytesIO(file.read()))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return '\n\n'.join(text_parts)
    
    elif ext == 'docx':
        from docx import Document
        from io import BytesIO
        doc = Document(BytesIO(file.read()))
        text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
        return '\n\n'.join(text_parts)
    
    elif ext == 'xlsx':
        import openpyxl
        from io import BytesIO
        wb = openpyxl.load_workbook(BytesIO(file.read()))
        text_parts = []
        for sheet in wb.worksheets:
            text_parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    text_parts.append(row_text)
        return '\n'.join(text_parts)
    
    elif ext == 'pdf':
        from PyPDF2 import PdfReader
        from io import BytesIO
        reader = PdfReader(BytesIO(file.read()))
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return '\n\n'.join(text_parts)
    
    return ""

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@study_bp.route('/study')
@login_required
def study_page():
    files = StudyFile.query.filter_by(user_id=current_user.id).all()
    subjects = Subject.query.filter_by(user_id=current_user.id).all()
    return render_template('study.html', files=files, subjects=subjects)

@study_bp.route('/chat/history/<int:file_id>')
@login_required
def get_chat_history(file_id):
    """Get chat history for a specific file"""
    conversations = BotConversation.query.filter_by(
        user_id=current_user.id,
        file_id=file_id
    ).order_by(BotConversation.sent_at.asc()).all()
    
    return jsonify([{
        'role': c.role,
        'content': c.content,
        'sent_at': c.sent_at.isoformat()
    } for c in conversations])

@study_bp.route('/chat/save', methods=['POST'])
@login_required
def save_chat_message():
    """Save a chat message"""
    data = request.json
    
    conversation = BotConversation(
        user_id=current_user.id,
        file_id=data.get('file_id'),
        role=data.get('role'),
        content=data.get('content')
    )
    db.session.add(conversation)
    db.session.commit()
    
    return jsonify({'success': True})

@study_bp.route('/chat/clear/<int:file_id>', methods=['POST'])
@login_required
def clear_chat_history(file_id):
    """Clear chat history for a file"""
    BotConversation.query.filter_by(
        user_id=current_user.id,
        file_id=file_id
    ).delete()
    db.session.commit()
    
    return jsonify({'success': True})

@study_bp.route('/upload', methods=['POST'])
@login_required
def upload_file():
    # Support both single file ('file') and multiple files ('files')
    files = request.files.getlist('files')
    
    # Filter out empty file entries
    files = [f for f in files if f and f.filename != '']
    
    if not files:
        # Fallback for single file upload (backward compatibility)
        single_file = request.files.get('file')
        if single_file and single_file.filename != '':
            files = [single_file]
    
    if not files:
        return jsonify({'error': 'No files provided'}), 400
    
    subject_id = request.form.get('subject_id')
    
    # Require subject selection
    if not subject_id:
        return jsonify({'error': 'Please select a subject'}), 400
    
    uploaded_files = []
    errors = []
    
    for file in files:
        if file.filename == '':
            continue
            
        if not allowed_file(file.filename):
            errors.append(f'{file.filename}: File type not allowed')
            continue
        
        filename = secure_filename(file.filename)
        
        try:
            content = extract_text_from_file(file, filename)
            
            if not content or len(content.strip()) < 10:
                errors.append(f'{file.filename}: Could not extract text or file is empty')
                continue
            
            study_file = StudyFile(
                user_id=current_user.id,
                subject_id=int(subject_id) if subject_id else None,
                filename=filename,
                original_name=file.filename,
                content=content
            )
            db.session.add(study_file)
            db.session.flush()  # Get the ID before commit
            
            # Get subject info for response
            subject_info = None
            if study_file.subject_id:
                subject = Subject.query.get(study_file.subject_id)
                if subject:
                    subject_info = {
                        'id': subject.id,
                        'name': subject.name,
                        'icon': subject.icon,
                        'color': subject.color,
                        'iconoir_icon': subject.iconoir_icon
                    }
            
            uploaded_files.append({
                'file_id': study_file.id,
                'filename': filename,
                'original_name': file.filename,
                'uploaded_at': datetime.utcnow().strftime('%b %d'),
                'subject': subject_info
            })
            
        except Exception as e:
            errors.append(f'{file.filename}: Error processing - {str(e)}')
    
    if uploaded_files:
        db.session.commit()
        
        response = {
            'success': True,
            'uploaded_count': len(uploaded_files),
            'files': uploaded_files
        }
        
        if errors:
            response['warnings'] = errors
        
        return jsonify(response)
    
    # No files were uploaded successfully
    error_msg = '; '.join(errors) if errors else 'No valid files to upload'
    return jsonify({'error': error_msg}), 400

from flask import session as flask_session

@study_bp.route('/bot/action', methods=['POST'])
@login_required
def bot_action():
    data = request.json
    file_id = data.get('file_id')
    action = data.get('action')  # 'quiz', 'question', 'ask'
    user_input = data.get('input', '')
    config = data.get('config', {})  # Quiz configuration: count, type
    
    study_file = StudyFile.query.get(file_id)
    if not study_file or study_file.user_id != current_user.id:
        return jsonify({'error': 'File not found'}), 404
    
    # Get or create conversation history for this file
    history_key = f'bot_history_{current_user.id}_{file_id}'
    conversation_history = flask_session.get(history_key, [])
    
    # Create bot with history
    bot = StudyBot(study_file.content, conversation_history)
    
    if action == 'quiz':
        # Extract quiz configuration parameters
        count = config.get('count', 5)
        question_type = config.get('type', 'mixed')
        result = bot.generate_quiz(num_questions=count, question_type=question_type)
    elif action == 'flashcards':
        result = bot.generate_flashcards()
    elif action == 'question':
        result = bot.ask_question()
    elif action == 'ask':
        result = bot.answer_question(user_input)
    elif action == 'check_answer':
        question = data.get('question', '')
        result = bot.check_answer(question, user_input)
    else:
        return jsonify({'error': 'Invalid action'}), 400
    
    # Save updated history back to session
    flask_session[history_key] = bot.get_history()
    
    return jsonify(result)

@study_bp.route('/bot/clear-memory/<int:file_id>', methods=['POST'])
@login_required
def clear_bot_memory(file_id):
    """Clear bot's conversation memory for a file"""
    history_key = f'bot_history_{current_user.id}_{file_id}'
    if history_key in flask_session:
        del flask_session[history_key]
    return jsonify({'success': True})

@study_bp.route('/track/quiz', methods=['POST'])
@login_required
def track_quiz():
    """Track quiz results for dashboard stats"""
    data = request.json
    file_id = data.get('file_id')
    total = data.get('total', 0)
    correct = data.get('correct', 0)
    
    study_file = StudyFile.query.get(file_id)
    topic = study_file.original_name if study_file else 'Quiz'
    
    # Create a study session record
    session = StudySession(
        user_id=current_user.id,
        topic=topic,
        duration_minutes=5,  # Estimate 5 mins per quiz
        questions_answered=total,
        correct_answers=correct,
        ended_at=datetime.utcnow()
    )
    db.session.add(session)
    
    # Update subject progress if file has a subject
    if study_file and study_file.subject_id:
        progress = SubjectProgress.query.filter_by(
            user_id=current_user.id,
            subject_id=study_file.subject_id
        ).first()
        
        if not progress:
            progress = SubjectProgress(
                user_id=current_user.id,
                subject_id=study_file.subject_id
            )
            db.session.add(progress)
        
        progress.questions_answered = (progress.questions_answered or 0) + total
        progress.correct_answers = (progress.correct_answers or 0) + correct
        progress.sessions_count = (progress.sessions_count or 0) + 1
        progress.study_minutes = (progress.study_minutes or 0) + 5
        progress.last_studied = datetime.utcnow()
    
    db.session.commit()
    
    return jsonify({'success': True})

@study_bp.route('/session/start', methods=['POST'])
@login_required
def start_session():
    data = request.json
    session = StudySession(
        user_id=current_user.id,
        topic=data.get('topic', 'General Study')
    )
    db.session.add(session)
    db.session.commit()
    return jsonify({'session_id': session.id})

@study_bp.route('/session/end', methods=['POST'])
@login_required
def end_session():
    data = request.json
    session = StudySession.query.get(data.get('session_id'))
    
    if session and session.user_id == current_user.id:
        session.ended_at = datetime.utcnow()
        session.duration_minutes = data.get('duration', 0)
        session.questions_answered = data.get('questions', 0)
        session.correct_answers = data.get('correct', 0)
        db.session.commit()
        return jsonify({'success': True})
    
    return jsonify({'error': 'Session not found'}), 404

@study_bp.route('/file/delete/<int:file_id>', methods=['POST'])
@login_required
def delete_file(file_id):
    """Delete a study file"""
    study_file = StudyFile.query.get(file_id)
    
    if not study_file or study_file.user_id != current_user.id:
        return jsonify({'error': 'File not found'}), 404
    
    # Delete associated chat history
    BotConversation.query.filter_by(
        user_id=current_user.id,
        file_id=file_id
    ).delete()
    
    # Clear bot memory for this file
    history_key = f'bot_history_{current_user.id}_{file_id}'
    if history_key in flask_session:
        del flask_session[history_key]
    
    # Delete the file
    db.session.delete(study_file)
    db.session.commit()
    
    return jsonify({'success': True})
